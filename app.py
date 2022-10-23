import streamlit as st
from pptx import Presentation

# BUILD GRAPHS
# from pptx.chart.data import CategoryChartData
# from pptx.enum.chart import XL_CHART_TYPE

# EDIT CHART
# from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

# from pptx.enum.dml import MSO_THEME_COLOR
# # Create shapes
# from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from io import BytesIO

st.set_page_config(page_title="AHGCC_LYRIC_CONVERTER", layout="wide")
# To hide hamburger (top right corner) and “Made with Streamlit” footer,
hide_streamlit_style = """
                       <style>
                       #MainMenu {visibility: hidden;}
                       footer {visibility: hidden;}
                       </style>
                       """
st.markdown(hide_streamlit_style, unsafe_allow_html=True)
st.markdown(
    "<h3 style='text-align: center; color: blue;'>AHGCC HURUMA LYRIC CONVERTER</h3>",
    unsafe_allow_html=True)

text=st.text_area(label="Paste your lyric here:")
# print(len(text.split("\n\n")))
lyric_title=text.split('\n')[0]

# print(type(text))
# if len(text.split("\n\n"))>1:
if text != "":

    pr2 = Presentation()

    #  Register the slide
    slide1_register = pr2.slide_layouts[0]  # 1 is title ONLY layout
    # add the second slide to presentation
    slide1 = pr2.slides.add_slide(slide1_register)

    title1 = slide1.placeholders[0]
    # Placeholder= Item in layout
    title1.text = text.split("\n")[0]

    title1 = slide1.placeholders[1]
    # Placeholder= Item in layout
    title1.text = "Lyric"

    for count, i in enumerate(text.split("\n\n")):
        # Register the slide
        slide_register = pr2.slide_layouts[6]  # 1 is title ONLY layout
        # add the second slide to presentation
        slide = pr2.slides.add_slide(slide_register)

        left = Inches(0.5)
        top = Inches(1)
        width = Inches(8.5)
        height = Inches(5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        p.text = i
        p.font.size = Pt(32)
    st.success(f"SUCCESSFULLY CONVERTED TO .ppt")
    # pr2.save("kigooco.pptx")

    # save the output into binary form
    binary_output = BytesIO()
    pr2.save(binary_output)

    st.download_button(label = f'Download {lyric_title}... in Powerpoint',
                       data = binary_output.getvalue(),
                       file_name = f'{lyric_title}.pptx')